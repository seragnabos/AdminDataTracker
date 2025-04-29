import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

#-------------------------------------
# دالة إنشاء الهيكل التنظيمي
#-------------------------------------
def create_org_chart(df, color_scheme="default", style="hierarchical", show_details=True):
    """إنشاء الهيكل التنظيمي التفاعلي مع خيارات تخصيص"""
    if 'الادارة' not in df.columns:
        return None

    departments = df['الادارة'].unique()
    fig = go.Figure()

    # إنشاء مستويات الهيكل التنظيمي
    y_positions = {}
    current_y = 0

    # معلومات إضافية لكل إدارة
    dept_info = {}
    for dept in departments:
        dept_data = df[df['الادارة'] == dept]
        dept_info[dept] = {
            'عدد الموظفين': len(dept_data),
            'الفئات الوظيفية': dept_data['فئة الوظيفة'].value_counts().to_dict() if 'فئة الوظيفة' in df.columns else {},
            'متوسط العمر': round((datetime.now() - pd.to_datetime(dept_data['تاريخ الميلاد'])).mean().days / 365.25, 1) if 'تاريخ الميلاد' in df.columns else None
        }

    # المستوى الأول - الإدارة العليا
    fig.add_trace(go.Scatter(
        x=[0], y=[current_y],
        mode='markers+text',
        text=['الإدارة العليا'],
        textposition='middle center',
        marker=dict(size=30, color='#1f77b4'),
        hoverinfo='text'
    ))

    # المستوى الثاني - الإدارات
    current_y -= 1
    x_position = -len(departments) / 2
    for dept in departments:
        y_positions[dept] = current_y
        fig.add_trace(go.Scatter(
            x=[x_position], y=[current_y],
            mode='markers+text',
            text=[dept],
            textposition='middle center',
            marker=dict(size=25, color='#2ca02c'),
            hoverinfo='text'
        ))

        # خطوط الربط مع الإدارة العليا
        fig.add_trace(go.Scatter(
            x=[0, x_position],
            y=[0, current_y],
            mode='lines',
            line=dict(color='#777', width=1),
            hoverinfo='none',
            showlegend=False
        ))

        x_position += 1

    # إعدادات المخطط
    fig.update_layout(
        showlegend=False,
        plot_bgcolor='white',
        title={
            'text': 'الهيكل التنظيمي',
            'x': 0.5,
            'xanchor': 'center',
            'font': dict(size=20, family='Tajawal')
        },
        font=dict(family='Tajawal'),
        height=600,
        xaxis=dict(showgrid=False, zeroline=False, showticklabels=False),
        yaxis=dict(showgrid=False, zeroline=False, showticklabels=False)
    )

    return fig

#-------------------------------------
# دالة إنشاء ملف باوربوينت
#-------------------------------------
def create_org_chart_pptx(df):
    """إنشاء ملف باوربوينت للهيكل التنظيمي"""
    prs = Presentation()

    # شريحة العنوان
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = "الهيكل التنظيمي"

    # شريحة الإحصائيات العامة
    stats_slide = prs.slides.add_slide(prs.slide_layouts[1])
    stats_slide.shapes.title.text = "إحصائيات الموظفين"
    stats_content = stats_slide.placeholders[1]

    dept_counts = df['الادارة'].value_counts()
    stats_text = f"""
    • إجمالي عدد الموظفين: {len(df)}
    • عدد الإدارات: {len(dept_counts)}
    • أكبر إدارة: {dept_counts.index[0]} ({dept_counts.iloc[0]} موظف)
    """
    stats_content.text = stats_text

    # شريحة لكل إدارة
    for dept in df['الادارة'].unique():
        dept_slide = prs.slides.add_slide(prs.slide_layouts[1])
        dept_slide.shapes.title.text = f"إدارة {dept}"
        dept_data = df[df['الادارة'] == dept]

        content = dept_slide.placeholders[1]
        text = f"""
        • عدد الموظفين: {len(dept_data)}
        • الفئات الوظيفية:
        """
        if 'فئة الوظيفة' in dept_data.columns:
            for job_cat, count in dept_data['فئة الوظيفة'].value_counts().items():
                text += f"\n  - {job_cat}: {count} موظف"

        content.text = text

    # حفظ الملف في الذاكرة
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

#-------------------------------------
# عرض التحليلات المتقدمة
#-------------------------------------
def display_advanced_analytics():
    """عرض التحليلات المتقدمة والهيكل التنظيمي"""
    if 'df' not in st.session_state or st.session_state.df is None:
        st.warning("لا توجد بيانات للتحليل")
        return

    df = st.session_state.df

    # إعدادات جانبية
    st.sidebar.markdown("### خيارات عرض الهيكل التنظيمي")

    def generate_unique_key(*parts) -> str:
        return "_".join(str(part) for part in parts)

    # استخدام generate_unique_key لكل مفتاح
    chart_style = st.sidebar.selectbox(
        "نمط عرض الهيكل التنظيمي",
        options=["هرمي", "شجري", "دائري"],
        key="org_chart_style_display"
    )

    color_scheme = st.sidebar.selectbox(
        "نظام الألوان المستخدم",
        options=["الافتراضي", "متدرج", "عالي التباين"],
        key="org_chart_color_scheme"
    )

    show_stats = st.sidebar.multiselect(
        "إظهار الإحصائيات التفصيلية",
        options=["عدد الموظفين", "متوسط العمر", "نسبة الذكور", "المؤهلات العلمية"],
        default=["عدد الموظفين"],
        key="org_chart_stats_display"
    )

    selected_dept = st.sidebar.multiselect(
        "تصفية حسب الإدارات",
        options=df['الادارة'].unique().tolist(),
        key="org_chart_dept_filter"
    )

    if selected_dept:
        df = df[df['الادارة'].isin(selected_dept)]

    st.markdown("## 🏢 الهيكل التنظيمي")

    # عرض المخطط
    org_chart = create_org_chart(df)
    if org_chart:
        st.plotly_chart(org_chart, use_container_width=True)

    # أزرار التصدير
    col1, col2 = st.columns(2)
    with col1:
        if st.button("تصدير كملف PowerPoint", key=generate_unique_key("export_pptx", "export_section", 1)):
            pptx_buffer = create_org_chart_pptx(df)
            st.download_button(
                label="تحميل العرض التقديمي",
                data=pptx_buffer,
                file_name="organizational_structure.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

    # إحصائيات عامة
    st.markdown("### 📊 إحصائيات الهيكل التنظيمي")

    col1, col2, col3 = st.columns(3)

    with col1:
        st.metric("إجمالي عدد الموظفين", len(df))

    with col2:
        st.metric("عدد الإدارات", df['الادارة'].nunique())

    with col3:
        if 'فئة الوظيفة' in df.columns:
            st.metric("عدد الفئات الوظيفية", df['فئة الوظيفة'].nunique())

    # تفاصيل الإدارات
    st.markdown("### 📋 تفاصيل الإدارات")
    for dept in df['الادارة'].unique():
        dept_data = df[df['الادارة'] == dept]
        with st.expander(f"إدارة {dept}"):
            col1, col2 = st.columns(2)

            with col1:
                st.metric("عدد الموظفين", len(dept_data))

                if 'فئة الوظيفة' in dept_data.columns:
                    job_cats = dept_data['فئة الوظيفة'].value_counts()
                    fig = px.pie(
                        values=job_cats.values,
                        names=job_cats.index,
                        title='توزيع الفئات الوظيفية'
                    )
                    st.plotly_chart(fig)

            with col2:
                if 'المؤهل العلمي' in dept_data.columns:
                    edu_counts = dept_data['المؤهل العلمي'].value_counts()
                    fig = px.bar(
                        x=edu_counts.index,
                        y=edu_counts.values,
                        title='المؤهلات العلمية'
                    )
                    st.plotly_chart(fig)